#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""WPS 演示助手 · 用 python-pptx 生成 .pptx(WPS 演示与 PowerPoint 通用格式)。

依赖:pip install python-pptx
既可被 Codex 当函数库调用,也可直接运行生成一份示例演示:
    python pptx_tools.py 产品介绍.pptx
"""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INK = RGBColor(0x11, 0x11, 0x11)
CORAL = RGBColor(0xD9, 0x77, 0x57)


def new_presentation():
    """新建 16:9 演示文稿。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, title, subtitle=""):
    """封面页:大标题 + 副标题。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_bullets_slide(prs, title, bullets):
    """要点页:标题 + 项目符号列表。bullets 可为 (文本, 层级) 或纯文本。"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(bullets):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(22 if level == 0 else 18)
    return slide


def add_section_slide(prs, text):
    """章节分隔页:居中大字 + 珊瑚色。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(1.5))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = CORAL
    return slide


def build_demo(path="示例-产品介绍.pptx"):
    """生成一份示例演示,验证依赖与流程。"""
    prs = new_presentation()
    add_title_slide(prs, "百灵鸟 EchoBird", "像专家一样部署 AI · WPS 办公三件套技能演示")
    add_bullets_slide(prs, "目录", ["背景与问题", "解决方案", "三大技能", "下一步"])
    add_section_slide(prs, "三大技能")
    add_bullets_slide(prs, "WPS 办公三件套", [
        "WPS 文档助手 — 一句话生成周报、合同、报告",
        ("基于 python-docx,WPS / Word 通用", 1),
        "WPS 表格助手 — 批量写数据、公式、分析",
        ("基于 openpyxl", 1),
        "WPS 演示助手 — 自动生成幻灯片",
        ("基于 python-pptx,即本页", 1),
    ])
    prs.save(path)
    return path


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "示例-产品介绍.pptx"
    saved = build_demo(out)
    print("已生成 WPS/PowerPoint 演示:", saved)
    p = Presentation(saved)
    print(f"--- 读回校验:共 {len(p.slides)} 页 ---")
